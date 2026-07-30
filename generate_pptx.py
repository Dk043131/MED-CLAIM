#!/usr/bin/env python3
"""
generate_pptx.py — Generates a 10-slide PowerPoint presentation (.pptx)
for MED-CLAIM (Knowledge Institute Of Technology)
"""

import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed yet.")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (MED-CLAIM Dark Theme)
    BG_COLOR = RGBColor(11, 15, 25)        # #0b0f19
    SURFACE_COLOR = RGBColor(17, 24, 39)   # #111827
    CARD_COLOR = RGBColor(31, 41, 55)      # #1f2937
    INDIGO = RGBColor(99, 102, 241)        # #6366f1
    CYAN = RGBColor(6, 182, 212)          # #06b6d4
    GREEN = RGBColor(16, 185, 129)        # #10b981
    WHITE = RGBColor(249, 250, 251)       # #f9fafb
    MUTED = RGBColor(156, 163, 175)       # #9ca3af

    blank_layout = prs.slide_layouts[6]

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.color.rgb = BG_COLOR
        return bg

    def add_header(slide, title_text, category_text="MED-CLAIM PRESENTATION"):
        # Header Badge / Category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.font.name = "Inter"

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.font.name = "Inter"

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "MED-CLAIM"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.font.name = "Inter"

    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Automated Medical Claim Adjudication Engine"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.font.name = "Inter"
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Ayushman Bharat (PM-JAY) 6-Gate Processing Pipeline & Multi-Language Engine"
    p3.font.size = Pt(15)
    p3.font.color.rgb = MUTED
    p3.font.name = "Inter"
    p3.space_before = Pt(12)

    # Presenters Card Box
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.5), Inches(11.333), Inches(2.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = SURFACE_COLOR
    card1.line.color.rgb = INDIGO

    p_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.7), Inches(11), Inches(1.8))
    tf_p = p_box.text_frame
    tf_p.word_wrap = True

    p_h = tf_p.paragraphs[0]
    p_h.text = "PROJECT PRESENTERS"
    p_h.font.size = Pt(12)
    p_h.font.bold = True
    p_h.font.color.rgb = CYAN
    p_h.font.name = "Inter"

    p_names = tf_p.add_paragraph()
    p_names.text = "Deepakkumar A  •  Dinakar S  •  Nandhakishore N  •  Devapriyan MJ"
    p_names.font.size = Pt(18)
    p_names.font.bold = True
    p_names.font.color.rgb = WHITE
    p_names.font.name = "Inter"
    p_names.space_before = Pt(8)

    p_inst = tf_p.add_paragraph()
    p_inst.text = "Knowledge Institute Of Technology (KIOT) — Department of CSE"
    p_inst.font.size = Pt(14)
    p_inst.font.color.rgb = MUTED
    p_inst.font.name = "Inter"
    p_inst.space_before = Pt(8)


    # Helper to add 4 grid cards
    def add_4_cards(slide, cards_data):
        positions = [
            (Inches(0.8), Inches(1.8)), (Inches(6.8), Inches(1.8)),
            (Inches(0.8), Inches(4.5)), (Inches(6.8), Inches(4.5))
        ]
        for idx, (title, desc) in enumerate(cards_data):
            x, y = positions[idx]
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.3))
            shape.fill.solid()
            shape.fill.fore_color.rgb = SURFACE_COLOR
            shape.line.color.rgb = CARD_COLOR

            tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.3), Inches(1.9))
            tf = tb.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            p1.text = title
            p1.font.size = Pt(18)
            p1.font.bold = True
            p1.font.color.rgb = CYAN
            p1.font.name = "Inter"

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(13)
            p2.font.color.rgb = MUTED
            p2.font.name = "Inter"
            p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_header(slide2, "Challenges in Medical Claim Adjudication", "PROBLEM STATEMENT")

    add_4_cards(slide2, [
        ("Manual Settlement Delays", "Insurance caseworkers spend 25–45 mins per bill verifying line items, causing claim clearance backlogs up to 30 days."),
        ("High Operational Errors", "Human coding errors in ICD-10/SNOMED CT mapping result in over $14B annually in wasted administrative expenses."),
        ("Fraud & Duplicate Billing", "Lack of real-time fingerprinting allows twin claims submitted across multiple clinics to bypass traditional audits."),
        ("Government Scheme Friction", "Complex 3-gate verification leads to accidental rejection of eligible Ayushman Bharat (PM-JAY) hospital claims.")
    ])

    # -------------------------------------------------------------
    # SLIDE 3: Proposed Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_header(slide3, "The MED-CLAIM Adjudication Platform", "PROPOSED SOLUTION")

    add_4_cards(slide3, [
        ("Autonomous 6-Gate Pipeline", "End-to-end processing pipeline evaluating OCR document intelligence, clinical coding, eligibility, and fraud risk in <400ms."),
        ("Human-in-the-Loop (HITL) Queue", "Intelligent triage automatically routes high-confidence claims (≥90%) to Auto-Approval while routing handwritten bills to caseworker review."),
        ("PM-JAY Ayushman Bharat Integration", "Built-in 3-gate eligibility checker, Aadhaar e-KYC card generator, and real-time family ₹5 Lakh annual cap tracking engine."),
        ("Multi-Format Intelligence", "Multi-engine OCR extracting clinical diagnosis, lab parameters, and line items from clean bills, lab reports, and doctor prescriptions.")
    ])

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_header(slide4, "6-Stage End-to-End Processing Architecture", "SYSTEM ARCHITECTURE")

    stages_data = [
        ("1. OCR / IDP", "Document intelligence for scanned bills & notes"),
        ("2. Structuring", "Extracts patient metadata, facility & line items"),
        ("3. ICD Coding", "Harmonizes diagnostic codes with ICD-10/SNOMED CT"),
        ("4. Eligibility", "Validates welfare policy & 7-day duplicate window"),
        ("5. Fraud Score", "Price anomaly guardrails & risk score evaluation"),
        ("6. Verdict", "Auto-Approval or escalation to PM-JAY portal")
    ]

    for idx, (s_title, s_desc) in enumerate(stages_data):
        x = Inches(0.8 + (idx % 3) * 3.9)
        y = Inches(1.8 if idx < 3 else 4.5)
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(2.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE_COLOR
        shape.line.color.rgb = INDIGO

        tb = slide4.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = s_title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = CYAN
        p1.font.name = "Inter"

        p2 = tf.add_paragraph()
        p2.text = s_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED
        p2.font.name = "Inter"
        p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 5: Multilingual & Mobile First
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_header(slide5, "Multilingual & Mobile-First Accessibility", "ACCESSIBILITY & I18N")

    add_4_cards(slide5, [
        ("10 Language Engine", "Full i18n support for 8 Indian regional languages (Hindi, Bengali, Telugu, Marathi, Tamil, Gujarati, Kannada, Malayalam, Punjabi) plus English and Spanish."),
        ("Instant Localization", "Dynamic DOM translation engine with persistent language selection via localStorage (`med_claim_lang`)."),
        ("Mobile Responsive Drawer", "Slide-over navigation drawer with backdrop overlay for screens ≤768px."),
        ("Touch Target Accessibility", "Minimum 44px touch height targets and scrollable table containers for smartphones.")
    ])

    # -------------------------------------------------------------
    # SLIDE 6: Document Intelligence
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_header(slide6, "Multi-Format Document Intelligence", "DOCUMENT INTELLIGENCE")

    add_4_cards(slide6, [
        ("Clean Hospital Bills", "Extracts itemized billing tables, admission/discharge dates, room rates, and tax calculations (>98% accuracy) for auto-approval."),
        ("Lab & Blood Reports", "Parses blood parameters (CBC, LFT, Widal, ESR), reference ranges, and abnormal value flags to validate clinical necessity."),
        ("Prescription Orders", "Reads prescribed medications, dosages, quantities, and doctor credentials to verify pharmacy claims."),
        ("Handwritten Notes", "Detects low OCR confidence (<70%) or ambiguous doctor handwriting, automatically flagging claims for caseworker HITL review.")
    ])

    # -------------------------------------------------------------
    # SLIDE 7: Fraud Detection
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_header(slide7, "Fraud Detection & Twin-Claim Prevention", "SECURITY GUARDRAILS")

    add_4_cards(slide7, [
        ("7-Day Twin Fingerprinting", "Fingerprints patient identity, provider ID, and diagnostic codes to instantly block duplicate claims submitted across clinics within a 7-day window."),
        ("Itemized Price Anomaly Scoring", "Evaluates unit prices against standardized NHA package rates to flag inflated billing items."),
        ("Soft vs Hard Escalation", "Fraud score 0.30–0.60 adds soft warning flags; score >0.60 triggers mandatory caseworker audit."),
        ("Audit Log Verification", "Generates immutable step-by-step reasoning panel documenting why every claim passed or failed.")
    ])

    # -------------------------------------------------------------
    # SLIDE 8: PM-JAY Integration
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    add_header(slide8, "Ayushman Bharat (PM-JAY) Scheme Management", "GOVERNMENT SCHEME PORTAL")

    add_4_cards(slide8, [
        ("3-Gate Eligibility Checker", "Verifies patient identity, hospital empanelment status, and 1,929 procedure package rates in real-time."),
        ("Aadhaar e-KYC Card Generator", "Simulates 12-digit Aadhaar OTP verification and generates downloadable PM-JAY Gold Ayushman Cards."),
        ("Family ₹5 Lakh Cap Tracker", "Monitors annual coverage balance, senior citizen separate caps, patient co-pay, and remaining family funds."),
        ("Portal Auto-Submission", "Automatically registers approved claims with NHA PM-JAY portal for instant disbursement processing.")
    ])

    # -------------------------------------------------------------
    # SLIDE 9: Metrics & Impact
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)
    add_header(slide9, "Performance Metrics & Benchmark Results", "RESULTS & IMPACT")

    metrics_data = [
        ("83.1%", "Auto-Adjudication Rate", "Claims auto-approved without manual intervention"),
        ("85%", "Staff Time Saved", "Reduction in caseworker processing effort"),
        ("<400ms", "Pipeline Latency", "End-to-end 6-gate execution speed"),
        ("94%", "Average Confidence", "OCR & clinical coding score accuracy")
    ]

    add_4_cards(slide9, [
        (f"{m[0]} — {m[1]}", m[2]) for m in metrics_data
    ])

    # -------------------------------------------------------------
    # SLIDE 10: Conclusion & Q&A
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10)

    title_box = slide10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf10 = title_box.text_frame
    tf10.word_wrap = True

    p = tf10.paragraphs[0]
    p.text = "Thank You! Any Questions?"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Inter"

    p2 = tf10.add_paragraph()
    p2.text = "MED-CLAIM sets a new benchmark in automated medical claim adjudication, bridging Indian healthcare schemes with state-of-the-art document intelligence."
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED
    p2.font.name = "Inter"
    p2.space_before = Pt(12)

    # Deployment links card
    card10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.4))
    card10.fill.solid()
    card10.fill.fore_color.rgb = SURFACE_COLOR
    card10.line.color.rgb = INDIGO

    p_box10 = slide10.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(11), Inches(2.0))
    tf_p10 = p_box10.text_frame
    tf_p10.word_wrap = True

    p_h10 = tf_p10.paragraphs[0]
    p_h10.text = "LIVE PROJECT DEPLOYMENT"
    p_h10.font.size = Pt(12)
    p_h10.font.bold = True
    p_h10.font.color.rgb = CYAN
    p_h10.font.name = "Inter"

    p_urls = tf_p10.add_paragraph()
    p_urls.text = "Live Demo: https://med-claim.vercel.app\nGitHub Repo: https://github.com/Dk043131/MED-CLAIM"
    p_urls.font.size = Pt(16)
    p_urls.font.bold = True
    p_urls.font.color.rgb = WHITE
    p_urls.font.name = "Inter"
    p_urls.space_before = Pt(8)

    p_inst10 = tf_p10.add_paragraph()
    p_inst10.text = "Knowledge Institute Of Technology (KIOT) • Department of CSE"
    p_inst10.font.size = Pt(14)
    p_inst10.font.color.rgb = MUTED
    p_inst10.font.name = "Inter"
    p_inst10.space_before = Pt(8)

    output_path = os.path.join(os.getcwd(), "med_claim_presentation.pptx")
    prs.save(output_path)
    print(f"SUCCESS: Created presentation at {output_path}")

if __name__ == "__main__":
    create_presentation()
